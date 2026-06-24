"""
Bank Churn – Step 4: 12-minute presentation deck (.pptx)
=========================================================
Builds an editable, presentation-quality PowerPoint (16:9, ~13 slides) following
the project rubric: opening · team · domain · KPIs · data prep · dashboard ·
deep-dive insights · recommendations · conclusion · questions.

Output: deliverables/Bank_Churn_Analysis.pptx
"""
from pathlib import Path
import json
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = Path(__file__).resolve().parent.parent
CH = BASE / "deliverables" / "charts"
K = json.loads((BASE / "deliverables" / "kpis.json").read_text())["kpis"]
df = pd.read_csv(BASE / "deliverables" / "bank_churn_clean.csv")

# ---- extra figures for the narrative ----------------------------------
ger = df[df.Geography == "Germany"]
ger_risk = ger[ger.Exited == 1].Balance.sum()
multi = df[df.NumOfProducts >= 3]
seg = df[(df.IsActiveMember == 0) & (df.NumOfProducts == 1)]
target_cut = 0.15
saved = int(round((K["churn_rate"] - target_cut) * K["total_customers"]))

# ---- palette & fonts --------------------------------------------------
NAVY = RGBColor(0x1F, 0x3A, 0x5F); TEAL = RGBColor(0x2A, 0x9D, 0x8F)
RED = RGBColor(0xE6, 0x39, 0x46); GREY = RGBColor(0x8D, 0x99, 0xAE)
DARK = RGBColor(0x2B, 0x2B, 0x2B); LIGHT = RGBColor(0xEE, 0xF2, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); AMBER = RGBColor(0xE9, 0xC4, 0x6A)
FONT = "Calibri"; FONT_H = "Calibri"
EMU = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ---- helpers ----------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, l, t, w, h, fill, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def rrect(s, l, t, w, h, fill):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.adjustments[0] = 0.06
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def txt(s, l, t, w, h, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, wrap=True):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf

def run(p, text, size, color=DARK, bold=False, font=FONT, italic=False):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = color
    return r

def set_bg(s, color):
    bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = color

def title_bar(s, title, kicker=None):
    rect(s, 0, 0, SW, 1.15, NAVY)
    rect(s, 0, 1.15, SW, 0.07, TEAL)
    if kicker:
        tf = txt(s, 0.55, 0.12, 12, 0.3)
        run(tf.paragraphs[0], kicker.upper(), 12, TEAL, bold=True)
        tf2 = txt(s, 0.55, 0.40, 12.2, 0.7)
        run(tf2.paragraphs[0], title, 27, WHITE, bold=True, font=FONT_H)
    else:
        tf = txt(s, 0.55, 0.18, 12.2, 0.85, anchor=MSO_ANCHOR.MIDDLE)
        run(tf.paragraphs[0], title, 29, WHITE, bold=True, font=FONT_H)

def footer(s, n):
    tf = txt(s, 0.55, 7.12, 8, 0.3)
    run(tf.paragraphs[0], "Bank Churn  ·  Data Analyst Program, Hebrew University of Jerusalem", 9, GREY)
    tf2 = txt(s, SW-1.4, 7.12, 0.9, 0.3, align=PP_ALIGN.RIGHT)
    run(tf2.paragraphs[0], str(n), 10, GREY, bold=True)

def bullets(s, items, l, t, w, h, size=15, gap=8):
    tf = txt(s, l, t, w, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.08
        if isinstance(it, tuple):
            lead, rest = it
            run(p, "▸ ", size, TEAL, bold=True)
            run(p, lead, size, NAVY, bold=True)
            run(p, rest, size, DARK)
        else:
            run(p, "▸ ", size, TEAL, bold=True)
            run(p, it, size, DARK)
    return tf

def img(s, path, box_l, box_t, box_w, box_h, align="center", valign="middle"):
    iw, ih = Image.open(path).size
    ar = iw / ih; bar = box_w / box_h
    if ar > bar:  # width-bound
        w = box_w; h = box_w / ar
    else:
        h = box_h; w = box_h * ar
    l = box_l + (box_w - w) * (0.5 if align == "center" else (1 if align == "right" else 0))
    t = box_t + (box_h - h) * (0.5 if valign == "middle" else (1 if valign == "bottom" else 0))
    s.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))

def kpi_card(s, l, t, w, h, value, label, sub, color, primary=False):
    card = rrect(s, l, t, w, h, NAVY if primary else WHITE)
    if not primary:
        card.line.color.rgb = LIGHT; card.line.width = Pt(1)
    rect(s, l, t, w, 0.09, color)  # top accent
    vc = WHITE if primary else color
    lc = RGBColor(0xCF,0xDA,0xE8) if primary else GREY
    sc = RGBColor(0xCF,0xDA,0xE8) if primary else GREY
    tf = txt(s, l+0.12, t+0.22, w-0.24, 0.55, align=PP_ALIGN.CENTER)
    run(tf.paragraphs[0], value, 26, vc, bold=True)
    tf2 = txt(s, l+0.1, t+0.78, w-0.2, 0.3, align=PP_ALIGN.CENTER)
    run(tf2.paragraphs[0], label.upper(), 10.5, (WHITE if primary else NAVY), bold=True)
    tf3 = txt(s, l+0.1, t+1.06, w-0.2, 0.4, align=PP_ALIGN.CENTER)
    run(tf3.paragraphs[0], sub, 9, sc)

def chip(s, l, t, w, text, color):
    c = rrect(s, l, t, w, 0.42, color)
    tf = txt(s, l, t+0.03, w, 0.36, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    run(tf.paragraphs[0], text, 12, WHITE, bold=True)

def callout(s, l, t, w, h, big, small, color=RED):
    rrect(s, l, t, w, h, LIGHT)
    rect(s, l, t, 0.10, h, color)
    tf = txt(s, l+0.28, t+0.12, w-0.4, h-0.2, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; run(p, big, 21, color, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(2); run(p2, small, 12.5, DARK)

# ===================================================================== 1 TITLE
s = slide(); set_bg(s, NAVY)
rect(s, 0, 0, 0.28, SH, TEAL)
rect(s, 0, 5.0, SW, 0.04, TEAL)
tf = txt(s, 1.0, 1.5, 11.3, 0.4)
run(tf.paragraphs[0], "BANK CHURN  ·  A BUSINESS PERSPECTIVE", 14, TEAL, bold=True)
tf = txt(s, 1.0, 2.0, 11.4, 1.9)
run(tf.paragraphs[0], "Why Are Our Customers Leaving?", 46, WHITE, bold=True, font=FONT_H)
tf = txt(s, 1.0, 3.7, 11.3, 0.7)
run(tf.paragraphs[0], "A churn & retention analysis of 10,000 retail-banking customers",
    19, RGBColor(0xD7,0xE0,0xEC))
tf = txt(s, 1.0, 5.25, 11.3, 0.6)
p = tf.paragraphs[0]
run(p, f"{K['churn_rate']:.1%} churn", 18, RED, bold=True)
run(p, "   ·   ", 18, GREY)
run(p, f"€{K['balance_at_risk']/1e6:.0f}M deposits at risk", 18, AMBER, bold=True)
run(p, "   ·   ", 18, GREY)
run(p, f"{K['churned_customers']:,} customers lost", 18, WHITE, bold=True)
tf = txt(s, 1.0, 6.5, 11.3, 0.5)
run(tf.paragraphs[0], "Data Analyst Program, Hebrew University of Jerusalem   |   David Zeff   |   June 2026",
    13, RGBColor(0xB6,0xC4,0xD6))

# ===================================================================== 2 TEAM
s = slide(); title_bar(s, "Team & Mentor", "Project")
rows = [("Presenter", "David Zeff"),
        ("Group members", "[ add teammate names ]"),
        ("Mentor", "[ add mentor name ]"),
        ("Dataset", "Bank Churn — 10,000 retail-banking customers"),
        ("Tools", "Python (pandas) · cleaning & analysis  |  Plotly · interactive dashboard  |  PowerPoint")]
y = 1.9
for lab, val in rows:
    rrect(s, 0.9, y, 2.7, 0.62, NAVY)
    tf = txt(s, 0.9, y+0.04, 2.7, 0.55, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    run(tf.paragraphs[0], lab, 14, WHITE, bold=True)
    tf = txt(s, 3.9, y+0.04, 8.4, 0.55, anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], val, 16, DARK)
    y += 0.82
footer(s, 2)

# ===================================================================== 3 DOMAIN
s = slide(); title_bar(s, "The Domain — Retail Banking Churn", "Introduction")
bullets(s, [
    ("Churn", " = a customer leaves the bank (closes their relationship). Our flag: Exited = 1."),
    ("Why it matters: ", "acquiring a new customer costs far more than keeping one — and a churned customer takes their deposits and fee income with them."),
    ("The base: ", "10,000 customers across France, Germany and Spain, with demographics, account balances, products held and activity."),
    ("Our questions: ", "How many leave? WHO leaves, and WHY? And WHERE should the bank intervene first?"),
], 0.7, 1.55, 7.4, 5.0, size=16, gap=14)
rrect(s, 8.5, 1.7, 4.2, 4.6, WHITE).line.color.rgb = LIGHT
img(s, CH/"01_overall_donut.png", 8.55, 1.8, 4.1, 4.4)
footer(s, 3)

# ===================================================================== 4 KPIs
s = slide(); title_bar(s, "KPIs & Success Metrics", "What we measure")
cards = [
    (f"{K['churn_rate']:.1%}", "Churn rate", "PRIMARY · target ≤ 15%", RED, True),
    (f"{1-K['churn_rate']:.1%}", "Retention rate", "keep > 85%", TEAL, False),
    (f"€{K['balance_at_risk']/1e6:.0f}M", "Balance at risk", f"{K['balance_at_risk_share']:.0%} of deposits", RED, False),
    (f"{K['active_member_rate']:.0%}", "Active-member rate", "engagement lever", NAVY, False),
    (f"{K['avg_products']:.2f}", "Products / customer", f"{K['single_product_share']:.0%} single-product", NAVY, False),
]
cw, gap = 2.34, 0.18; x = 0.62
for v, l, sub, c, prim in cards:
    kpi_card(s, x, 1.65, cw, 1.6, v, l, sub, c, prim); x += cw + gap
tf = txt(s, 0.62, 3.7, 12.1, 2.8)
run(tf.paragraphs[0], "Thought process", 16, NAVY, bold=True)
bullets(s, [
    ("Churn rate is the north-star metric ", "— the single number that tells us whether we are winning or losing customers."),
    ("Balance at risk", " translates churn into euros, so the business feels the impact (€186M ≈ a quarter of all deposits)."),
    ("Active-member rate & products-per-customer", " are the levers we can actually pull — engagement and product fit drive retention."),
    ("Success", " = move churn from 20% → ≤15% and lift retention above 85% over the next year."),
], 0.62, 4.15, 12.1, 2.6, size=14.5, gap=9)
footer(s, 4)

# ===================================================================== 5 DATA PREP
s = slide(); title_bar(s, "Data Preparation & Assumptions", "From messy to trusted")
bullets(s, [
    ("Source: ", "one Excel workbook, two sheets (Customer_Info + Account_Info), joined on CustomerId."),
    ("Standardised geography ", "— France was split across 'France/French/FRA' → unified."),
    ("Fixed types ", "— stripped '€' from Balance & EstimatedSalary (text→number); Yes/No → 1/0."),
    ("Removed duplicates", ", imputed 3 missing ages (median), reconciled a duplicated Tenure column."),
    ("Filters/assumptions: ", "snapshot data (no time dimension); full customer base analysed, no sampling."),
], 0.7, 1.5, 7.2, 5.0, size=14.5, gap=11)
# right column: data-quality fix + validation
callout(s, 8.1, 1.6, 4.6, 1.55,
        "Data-quality catch",
        "'HasCrCard' was identical to 'IsActiveMember' for every row — a corrupted column. "
        "Detected and restored to authentic values.", RED)
rrect(s, 8.1, 3.4, 4.6, 2.9, LIGHT)
tf = txt(s, 8.35, 3.55, 4.2, 0.4); run(tf.paragraphs[0], "Validation vs. published reference", 14, NAVY, bold=True)
checks = [("Rows", "10,000 ✓"), ("Missing values", "0 ✓"), ("Churn rate", "20.37% ✓"),
          ("Geography mix", "match ✓"), ("Total balance", "€0.77bn ✓")]
yy = 4.05
for a, b in checks:
    tf = txt(s, 8.35, yy, 2.6, 0.35); run(tf.paragraphs[0], a, 13, DARK)
    tf = txt(s, 10.7, yy, 1.9, 0.35, align=PP_ALIGN.RIGHT); run(tf.paragraphs[0], b, 13, TEAL, bold=True)
    yy += 0.44
footer(s, 5)

# ===================================================================== 6 DASHBOARD
s = slide(); title_bar(s, "The Dashboard — One Page, Six Lenses", "BI tool: Plotly")
rrect(s, 0.5, 1.45, 12.35, 5.35, WHITE).line.color.rgb = LIGHT
img(s, CH/"dashboard_full.png", 0.65, 1.55, 12.05, 5.15)
footer(s, 6)

# ===================================================================== 7 INSIGHT 1 PRODUCTS
s = slide(); title_bar(s, "Insight 1 — The Product Paradox", "Deep dive")
rrect(s, 0.55, 1.5, 6.5, 4.7, WHITE).line.color.rgb = LIGHT
img(s, CH/"03_products.png", 0.65, 1.65, 6.3, 4.4)
bullets(s, [
    ("2 products is the sweet spot ", "— just 7.6% churn. These customers are engaged and sticky."),
    ("3–4 products = near-certain exit ", "— 83% and 100% churn. The 326 customers here are almost all gone."),
    ("Half the base (51%) holds a single product ", "and churns at 27.7%."),
    ("Read: ", "piling products onto unhappy customers backfires; the goal is the right second product, not more products."),
], 7.3, 1.7, 5.4, 4.6, size=14.5, gap=12)
callout(s, 7.3, 5.55, 5.4, 0.9, "1 → 2 products",
        "moving single-product customers to two is the biggest retention prize.", TEAL)
footer(s, 7)

# ===================================================================== 8 INSIGHT 2 GERMANY
s = slide(); title_bar(s, "Insight 2 — Germany & the Deposit Drain", "Deep dive")
rrect(s, 0.5, 1.5, 6.15, 4.0, WHITE).line.color.rgb = LIGHT
img(s, CH/"02_geography.png", 0.6, 1.6, 5.95, 3.8)
rrect(s, 6.75, 1.5, 6.05, 4.0, WHITE).line.color.rgb = LIGHT
img(s, CH/"06_balance_at_risk.png", 6.85, 1.6, 5.85, 3.8)
callout(s, 0.5, 5.7, 12.3, 1.15,
        f"Germany: {ger.Exited.mean():.0%} churn · €{ger.Balance.mean()/1000:.0f}k avg balance · €{ger_risk/1e6:.0f}M walking out the door",
        "Germany churns at twice the rate of France & Spain, yet holds ~2× the average balance — "
        "so it is both our biggest leak and our most valuable one. Priority market #1.", RED)
footer(s, 8)

# ===================================================================== 9 INSIGHT 3 AGE/ENGAGEMENT
s = slide(); title_bar(s, "Insight 3 — Age & Engagement", "Deep dive")
rrect(s, 0.5, 1.5, 6.15, 3.9, WHITE).line.color.rgb = LIGHT
img(s, CH/"04_age.png", 0.6, 1.6, 5.95, 3.7)
rrect(s, 6.75, 1.5, 6.05, 3.9, WHITE).line.color.rgb = LIGHT
img(s, CH/"05_activity_gender.png", 6.8, 1.7, 5.95, 3.5)
callout(s, 0.5, 5.6, 12.3, 1.2,
        f"{len(seg):,} customers are BOTH inactive and single-product — and {seg.Exited.mean():.0%} of them churn",
        "Age 51–60 churns at 56%; inactive members at 27% (vs 14% active); women at 25% (vs 17%). "
        "The inactive + single-product overlap is the single most actionable at-risk segment.", RED)
footer(s, 9)

# ===================================================================== 10 WHAT DOESN'T
s = slide(); title_bar(s, "Equally Useful — What Does NOT Drive Churn", "Deep dive")
tf = txt(s, 0.7, 1.45, 12, 0.6)
run(tf.paragraphs[0], "Three things the bank might expect to matter — but the data says don't:", 16, DARK)
items = [("Tenure", "≈ 19–21% churn at every tenure", "Loyalty length is flat — long-tenured customers leave just as often."),
         ("Credit score", "22% (poor) vs 19% (good)", "Only a faint signal — not a useful targeting variable."),
         ("Has credit card", "20.8% vs 20.2%", "No effect. (This is also where we caught the corrupted column.)")]
x = 0.7
for title_, stat, desc in items:
    rrect(s, x, 2.3, 3.9, 3.4, WHITE).line.color.rgb = LIGHT
    rect(s, x, 2.3, 3.9, 0.09, GREY)
    tf = txt(s, x+0.25, 2.55, 3.4, 0.5); run(tf.paragraphs[0], title_, 18, NAVY, bold=True)
    tf = txt(s, x+0.25, 3.15, 3.4, 0.6); run(tf.paragraphs[0], stat, 16, RED, bold=True)
    tf = txt(s, x+0.25, 3.95, 3.45, 1.6); run(tf.paragraphs[0], desc, 13.5, DARK)
    x += 4.1
callout(s, 0.7, 5.95, 11.9, 0.8, "Focus the retention budget where it counts",
        "products, geography, age and engagement — not tenure, credit score or card ownership.", TEAL)
footer(s, 10)

# ===================================================================== 11 RECOMMENDATIONS
s = slide(); title_bar(s, "Recommendations", "From insight to action")
recs = [
    ("Launch a Germany retention task force", "32% churn + high balances → review pricing, product fit and local service. Biggest single opportunity."),
    ("Fix the multi-product problem", "Audit why 3–4-product customers leave (fees? mis-selling?) before any further cross-sell to them."),
    ("Move single-product customers to two", "Targeted, value-led second product — 2 products is the 7.6%-churn sweet spot. Largest reachable group."),
    ("Re-activate the inactive", f"Engagement campaign for the {len(seg):,} inactive single-product customers churning at {seg.Exited.mean():.0%}."),
    ("Protect the 45–60 segment", "Proactive relationship outreach for older, higher-churn, higher-balance customers (and women)."),
]
y = 1.55
for i, (head, body) in enumerate(recs, 1):
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y), Inches(0.62), Inches(0.62))
    circ.fill.solid(); circ.fill.fore_color.rgb = TEAL if i in (1,3,4) else NAVY; circ.line.fill.background()
    circ.shadow.inherit = False
    tf = circ.text_frame; tf.word_wrap=False; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; run(p, str(i), 20, WHITE, bold=True)
    tf = txt(s, 1.55, y-0.05, 11.0, 0.5); run(tf.paragraphs[0], head, 16.5, NAVY, bold=True)
    tf = txt(s, 1.55, y+0.40, 11.1, 0.6); run(tf.paragraphs[0], body, 13, DARK)
    y += 1.06
footer(s, 11)

# ===================================================================== 12 CONCLUSION
s = slide(); set_bg(s, NAVY); rect(s, 0, 0, 0.28, SH, TEAL)
tf = txt(s, 1.0, 0.9, 11, 0.4); run(tf.paragraphs[0], "CONCLUSION", 14, TEAL, bold=True)
tf = txt(s, 1.0, 1.4, 11.4, 1.7)
run(tf.paragraphs[0], "1 in 5 customers leave — but the churn is concentrated and addressable.",
    30, WHITE, bold=True, font=FONT_H)
bullets_white = [
    "Churn is driven by product fit, geography, age and engagement — not by tenure or credit history.",
    "Germany, multi-product and inactive single-product customers are the clear priority segments.",
    f"Cutting churn from {K['churn_rate']:.1%} to {target_cut:.0%} would retain ~{saved:,} more customers each cycle and protect tens of €M in deposits.",
]
tf = txt(s, 1.0, 3.4, 11.3, 2.4)
for i, b in enumerate(bullets_white):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(12); p.line_spacing = 1.1
    run(p, "▸  ", 17, TEAL, bold=True); run(p, b, 17, RGBColor(0xE3,0xEA,0xF2))
chip(s, 1.0, 6.25, 5.6, f"Target: churn ≤ {target_cut:.0%}  ·  retention ≥ 85%", TEAL)

# ===================================================================== 13 QUESTIONS
s = slide(); set_bg(s, NAVY)
rect(s, 0, 3.0, SW, 0.05, TEAL)
tf = txt(s, 0, 2.0, SW, 1.4, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
run(tf.paragraphs[0], "Questions?", 54, WHITE, bold=True, font=FONT_H)
tf = txt(s, 0, 3.3, SW, 0.6, align=PP_ALIGN.CENTER)
run(tf.paragraphs[0], "Thank you  ·  Bank Churn — A Business Perspective", 18, RGBColor(0xD7,0xE0,0xEC))
tf = txt(s, 0, 4.0, SW, 0.5, align=PP_ALIGN.CENTER)
run(tf.paragraphs[0], "David Zeff   ·   Data Analyst Program, Hebrew University of Jerusalem", 13, GREY)

out = BASE / "deliverables" / "Bank_Churn_Analysis.pptx"
prs.save(out)
print(f"Deck → {out}  ({len(prs.slides._sldIdLst)} slides)")
